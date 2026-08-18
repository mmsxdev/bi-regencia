# -*- coding: utf-8 -*-
"""
Gera o PowerPoint de apresentacao do BI de Regência (SENAI Vila Canaã).

Uso:
    python gerar_ppt.py

Requisitos: python-pptx, Pillow
Saida: BI_Regencia_Apresentacao.pptx
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
GRAF = os.path.join(HERE, "graficos")
OUT_PPTX = os.path.join(HERE, "BI_Regencia_Apresentacao.pptx")

# ---- Cores SENAI (iguais ao app) ----
BG = RGBColor(0x0B, 0x0F, 0x14)
CARD = RGBColor(0x11, 0x18, 0x27)
SURFACE = RGBColor(0x1A, 0x20, 0x29)
RED = RGBColor(0xD7, 0x19, 0x20)
AZUL = RGBColor(0x2E, 0x90, 0xFA)
TXT = RGBColor(0xF3, 0xF4, 0xF6)
SEC = RGBColor(0x9C, 0xA3, 0xAF)
TER = RGBColor(0x6B, 0x72, 0x80)
BORDA = RGBColor(0x2A, 0x30, 0x3A)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, pars, anchor=MSO_ANCHOR.TOP, wrap=True):
    """pars: list of (runs, space_after, align) ; runs: list of (text, size, bold, color, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    for i, (runs, space_after, align) in enumerate(pars):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for (t, size, bold, color, *rest) in runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Inter"
            if rest and rest[0]:
                r.font.italic = True
    return tb


def title_bar(s, title, subtitle=None):
    rect(s, Inches(0), Inches(0), SW, Inches(1.05), fill=CARD)
    rect(s, Inches(0), Inches(0), Inches(0.14), Inches(1.05), fill=RED)
    text(s, Inches(0.45), Inches(0.12), Inches(10.6), Inches(0.55),
         [([(title, 28, True, TXT)], 0, PP_ALIGN.LEFT)])
    if subtitle:
        text(s, Inches(0.45), Inches(0.62), Inches(11.9), Inches(0.4),
             [([(subtitle, 13, False, SEC)], 0, PP_ALIGN.LEFT)])
    text(s, Inches(11.6), Inches(0.32), Inches(1.5), Inches(0.45),
         [([("SENAI", 17, True, RED)], 0, PP_ALIGN.RIGHT)])


def bullets(s, x, y, w, h, items, size=15, gap=8, color=TXT):
    pars = []
    for it in items:
        if isinstance(it, tuple):
            txt, bold, col, lvl = it
        else:
            txt, bold, col, lvl = it, False, color, 0
        bullet = "   " * lvl + ("•  " if lvl == 0 else "–  ")
        runs = [(bullet, size, True, RED if lvl == 0 else AZUL)]
        runs.append((txt, size, bold, col))
        pars.append((runs, gap, PP_ALIGN.LEFT))
    return text(s, x, y, w, h, pars)


def kpi_card(s, x, y, w, h, label, value, hint, dot_color=None):
    rect(s, x, y, w, h, fill=CARD, line=BORDA)
    rect(s, x, y, w, Inches(0.07), fill=RED)
    # value
    t = text(s, x + Inches(0.25), y + Inches(0.28), w - Inches(0.5), Inches(0.9),
             [([(value, 34, True, TXT)], 0, PP_ALIGN.LEFT)], anchor=MSO_ANCHOR.MIDDLE)
    if dot_color:
        d = rect(s, x + Inches(0.27), y + Inches(0.46), Inches(0.18), Inches(0.18), fill=dot_color)
        d.rotation = 0
    # label
    text(s, x + Inches(0.25), y + Inches(1.05), w - Inches(0.5), Inches(0.35),
         [([(label, 13, True, SEC)], 0, PP_ALIGN.LEFT)])
    if hint:
        text(s, x + Inches(0.25), y + Inches(1.42), w - Inches(0.5), Inches(0.5),
             [([(hint, 11, False, TER)], 0, PP_ALIGN.LEFT)])


def pic_fit(s, path, x, y, max_w, max_h, border=True):
    im = Image.open(path)
    ratio = im.width / im.height
    w = max_w
    h = Emu(int(w * im.height / im.width))
    if h > max_h:
        h = max_h
        w = Emu(int(h * im.width / im.height))
    if border:
        rect(s, x, y, w, h, fill=None, line=BORDA)
    s.shapes.add_picture(path, x, y, w, h)
    return w, h


def note(s, txt, x=Inches(0.45), y=Inches(6.95), w=Inches(12.4), h=Inches(0.45)):
    text(s, x, y, w, h, [([(txt, 11.5, False, TER, True)], 0, PP_ALIGN.LEFT)])


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------
def s_capa():
    s = add_slide()
    rect(s, Inches(0), Inches(0), SW, SH, fill=BG)
    rect(s, Inches(0), Inches(2.55), SW, Inches(0.06), fill=BORDA)
    text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5),
         [([("SENAI · SERVIÇO NACIONAL DE APRENDIZAGEM INDUSTRIAL", 16, True, SEC)], 0, PP_ALIGN.LEFT)])
    text(s, Inches(0.8), Inches(2.78), Inches(11.7), Inches(1.3),
         [([("BI DE REGÊNCIA", 52, True, TXT)], 0, PP_ALIGN.LEFT)])
    text(s, Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.6),
         [([("Frequência dos instrutores em sala de aula — quadro 2026", 22, False, AZUL)], 0, PP_ALIGN.LEFT)])
    rect(s, Inches(0.8), Inches(4.4), Inches(4.4), Inches(0.045), fill=RED)
    text(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(0.5),
         [([("Escola SENAI Vila Canaã · Departamento Regional de Goiás", 15, False, SEC)], 0, PP_ALIGN.LEFT)])
    text(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.4),
         [([("Relatório de dados consolidados da planilha de regência", 13, False, TER)], 0, PP_ALIGN.LEFT)])


def s_objetivo():
    s = add_slide()
    title_bar(s, "O que este painel entrega", "Objetivo e proposito do BI")
    bullets(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(3.2), [
        "Consolidar as horas-aula (regência) de todos os instrutores do quadro em um só painel.",
        "Responder 4 perguntas de gestão: quem ministra? quanto? cumpre a carga prevista? onde está o risco?",
        "Transformar a planilha de lançamentos, que já existe, em informação de gestão.",
        "Ficar sempre atualizado: conexão com a planilha no SharePoint/OneDrive + botão de atualização manual.",
    ], size=17, gap=14)
    rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.045), fill=BORDA)
    text(s, Inches(0.55), Inches(4.8), Inches(12.2), Inches(0.45),
         [([("Conceito-chave — Frequência:", 16, True, TXT)], 0, PP_ALIGN.LEFT)])
    bullets(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.6), [
        "Frequência (%) = horas-aula realizadas ÷ carga de regência esperada no mês.",
        "100% = cumprida a carga. Acima de 100% = ministrou acima do previsto. Abaixo de 50% = alerta.",
    ], size=15, gap=8)


def s_fonte():
    s = add_slide()
    title_bar(s, "Fonte dos dados", "De onde vem a informação")
    bullets(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(3.0), [
        "Aba \"CONSOLIDADO \" da planilha REGÊNCIA - INSTRUTORES DO QUADRO 2026.xlsx.",
        "1 linha por instrutor com: DOCENTE, carga horária (Ch), ÁREA, H/AULA e % por mês, ANO, EXTRA-QUADRO.",
        "O painel lê a planilha automaticamente e organiza em: Visão Geral, Por Instrutor, Por Mês e Tabela.",
        "Atualização automática a cada 10 minutos + botão \"Atualizar dados agora\".",
    ], size=16, gap=12)
    rect(s, Inches(0.55), Inches(3.55), Inches(12.2), Inches(0.045), fill=BORDA)
    text(s, Inches(0.55), Inches(3.8), Inches(12.2), Inches(0.4),
         [([("Atenção na leitura:", 15, True, TXT)], 0, PP_ALIGN.LEFT)])
    bullets(s, Inches(0.55), Inches(4.3), Inches(12.2), Inches(2.4), [
        "O n.º de instrutores com lançamento cresce ao longo do ano (37 em jan → 45 em ago): entradas em 2026.",
        "Existem dois totais no painel: soma dos meses (37.306 h) e coluna ANO (36.458 h) — conferir padrão na planilha.",
    ], size=15, gap=9)


def s_kpis():
    s = add_slide()
    title_bar(s, "Números-chave do quadro (2026)", "Visão consolidada")
    kpi_card(s, Inches(0.55), Inches(1.5), Inches(2.9), Inches(2.3),
             "INSTRUTORES NO QUADRO", "45", "Cobertura completa desde agosto")
    kpi_card(s, Inches(3.65), Inches(1.5), Inches(2.9), Inches(2.3),
             "FREQUÊNCIA MÉDIA", "61%", "média das frequências mensais", RGBColor(0xA3, 0xE6, 0x35))
    kpi_card(s, Inches(6.75), Inches(1.5), Inches(2.9), Inches(2.3),
             "HORAS-AULA REALIZADAS", "37,3 mil", "soma dos 12 meses do período")
    kpi_card(s, Inches(9.85), Inches(1.5), Inches(2.9), Inches(2.3),
             "ABAIXO DE 50%", "12", "27% do quadro em zona de alerta", RED)
    text(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(0.5),
         [([("Leitura: a média de 61% esconde realidades distintas — parte do quadro cumpre bem a carga, outra parte está em zona crítica.",
             15, False, SEC)], 0, PP_ALIGN.LEFT)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(1.8), [
        "16 instrutores entre 70% e 100% de frequência.",
        "2 instrutores acima de 100% (ministram acima da carga prevista).",
        "15 instrutores entre 50% e 70% — área de atenção/monitoramento.",
    ], size=14, gap=8)


def s_sazonalidade():
    s = add_slide()
    title_bar(s, "Sazonalidade — frequência por mês", "Cair em julho não é risco, é calendário")
    pic_fit(s, os.path.join(GRAF, "mes_frequencia_line.png"),
            Inches(0.5), Inches(1.35), Inches(7.3), Inches(4.9))
    bullets(s, Inches(8.1), Inches(1.5), Inches(4.7), Inches(5.0), [
        "Janeiro: 41% — início de ano / entrada de contratos.",
        "Julho: 18% — férias escolares (recesso).",
        "Dezembro: 37% — encerramento do ano.",
        "Março a Setembro: 78–82% — período produtivo.",
        ("Meta única anual seria injusta: usar metas por período.", True, AZUL, 0),
    ], size=14, gap=10)


def s_volume():
    s = add_slide()
    title_bar(s, "Volume de horas-aula por mês", "Quando o quadro está mais ativo")
    pic_fit(s, os.path.join(GRAF, "mes_horas_bar.png"),
            Inches(0.5), Inches(1.35), Inches(7.3), Inches(4.9))
    bullets(s, Inches(8.1), Inches(1.5), Inches(4.7), Inches(5.0), [
        "Pico: agosto (4.403 h) e setembro (4.291 h).",
        "Vale: julho (890 h), janeiro (1.757 h), dezembro (1.990 h).",
        "Total do ano: ~37.306 horas-aula.",
        "Essa curva orienta a alocação de carga e a necessidade de reforço do quadro.",
    ], size=14, gap=10)


def s_distribuicao():
    s = add_slide()
    title_bar(s, "Distribuição da frequência média", "Onde o quadro se concentra")
    pic_fit(s, os.path.join(GRAF, "distribuicao_hist.png"),
            Inches(0.5), Inches(1.35), Inches(7.0), Inches(4.9))
    bullets(s, Inches(7.9), Inches(1.55), Inches(5.0), Inches(5.0), [
        "12 instrutores abaixo de 50% (alerta).",
        "15 entre 50% e 70% (monitorar).",
        "16 entre 70% e 100% (saudável).",
        "2 acima de 100% (referência de desempenho).",
        ("27% do quadro cumpre menos da metade da carga prevista.",
         True, RED, 0),
    ], size=14, gap=11)


def s_area():
    s = add_slide()
    title_bar(s, "Frequência média por área", "Risco é concentrado, não está espalhado")
    pic_fit(s, os.path.join(GRAF, "area_freq_bar.png"),
            Inches(0.5), Inches(1.35), Inches(7.6), Inches(5.0))
    bullets(s, Inches(8.35), Inches(1.55), Inches(4.5), Inches(5.0), [
        "Gráfica editorial (11 instr.): 74,7%.",
        "Alimentos e bebidas (6 instr.): 73,0%.",
        "TI (5 instr.): 69,5%.",
        "Manutenção automotiva (13 instr.): ~46%.",
        "Administração (2 instr.): 44,1%.",
        ("Área crítica: Manutenção automotiva.", True, RED, 0),
    ], size=14, gap=10)


def s_atencoes():
    s = add_slide()
    title_bar(s, "Pontos de atenção — menores frequências", "Identificados por nome e área")
    pic_fit(s, os.path.join(GRAF, "instrutores_menor.png"),
            Inches(0.5), Inches(1.35), Inches(7.6), Inches(5.0))
    bullets(s, Inches(8.35), Inches(1.55), Inches(4.5), Inches(5.0), [
        "Os 5 menores índices do quadro estão em Manutenção automotiva.",
        "DIOGO DE SOUZA PIMENTEL: 10,8% (207 h no ano).",
        "DINAIRON DA SILVA BORGES: 14,6% (280 h).",
        "CAIO CEZAR BRAZ E SILVA: 22,1% (383 h).",
        ("Pergunta à coordenação: carga reduzida? afastamento? realocação?"),
        ("O painel localiza quem e onde — a causa é decidida com a coordenação.", True, AZUL, 0),
    ], size=14, gap=10)


def s_automotiva():
    s = add_slide()
    title_bar(s, "Manutenção automotiva — área única", "O painel unifica; o polo vira filtro")
    pic_fit(s, os.path.join(GRAF, "automotiva_polos.png"),
            Inches(0.5), Inches(1.35), Inches(7.4), Inches(5.2))
    bullets(s, Inches(8.15), Inches(1.45), Inches(4.7), Inches(5.4), [
        "As \"Manutenção automotiva\" e \"(JD)\" foram unificadas em UMA área: 13 instrutores.",
        "(JD) = rótulo interno do polo da escola parceira (Col. Est. Jardim Vila Boa / SEDUC) — confirmar com a coordenação.",
        "Frequência média da área: ~46% — 7 de 13 abaixo de 50%.",
        ("Novo filtro POLO/LOCAL: separa Vila Canaã × Jardim Vila Boa (SEDUC) sem fragmentar a área.", True, AZUL, 0),
        "O painel lê a coluna POLO/LOCAL da planilha, se existir; senão, infere pelo rótulo antigo.",
        "Recomendação: preencher POLO/LOCAL para todos os instrutores na planilha-fonte.",
    ], size=13, gap=9)


def s_destaques():
    s = add_slide()
    title_bar(s, "Destaques do quadro", "Referências internas de alta frequência")
    pic_fit(s, os.path.join(GRAF, "instrutores_maior.png"),
            Inches(0.5), Inches(1.35), Inches(7.6), Inches(5.0))
    bullets(s, Inches(8.35), Inches(1.55), Inches(4.5), Inches(5.0), [
        "ROMULO FLORIANO LIMEIRA: 109,5% — maior frequência (1.314 h).",
        "BRUNA ARIEL DIAS GUARIGLIA: 103,8% (1.495 h).",
        "LEIDINA LAIS: 1.552 h no ano — maior volume, na própria área crítica.",
        ("Usar esses casos como referência de planejamento de carga.", True, AZUL, 0),
    ], size=14, gap=10)


def s_heatmap():
    s = add_slide()
    title_bar(s, "Visão 360° — heatmap instrutor × mês", "Frequência célula a célula")
    pic_fit(s, os.path.join(GRAF, "heatmap.png"),
            Inches(0.5), Inches(1.35), Inches(7.6), Inches(5.4))
    bullets(s, Inches(8.35), Inches(1.55), Inches(4.5), Inches(5.4), [
        "Cada célula = frequência de um instrutor em um mês.",
        "Colunas fracas em JAN/JUL/DEZ: sazonalidade.",
        "Linhas persistentemente vermelhas: caso de gestão.",
        "Células vazias = mês sem lançamento (conferir preenchimento).",
    ], size=14, gap=10)


def s_qualidade():
    s = add_slide()
    title_bar(s, "Qualidade dos dados — o BI também audita", "Falhas de preenchimento identificadas")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(4.2), [
        "Nomes de área com erros de digitação: \"Contrução Civil\", \"Grafica editorial\", e o rótulo \"(JD)\".",
        "Instrutores com meses sem lançamento: CHRISTIAN TEILOR (7 meses), JOÃO VICTOR MARTINS (5 meses).",
        "1 instrutor sem total anual preenchido (THAUANA MACHADO).",
        "Coluna EXTRA-QUADRO zerada no quadro (ainda não utilizada).",
        "Dois totais divergentes: soma dos meses (37.306 h) vs coluna ANO (36.458 h).",
    ], size=16, gap=12)
    rect(s, Inches(0.6), Inches(5.1), Inches(12.2), Inches(0.045), fill=BORDA)
    text(s, Inches(0.6), Inches(5.35), Inches(12.2), Inches(0.5),
         [([("Mensagem: corrigir a planilha-fonte é pré-requisito para os próximos ciclos do BI.", 16, True, AZUL)],
           0, PP_ALIGN.LEFT)])


def s_plano():
    s = add_slide()
    title_bar(s, "Plano de ação proposto", "Fechamento com decisões")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(4.6), [
        ("1. Validar com a coordenação os casos de baixa frequência (quem, por quê, o que fazer).", True, TXT, 0),
        ("    · Prioridade: Manutenção automotiva (7 de 13 abaixo de 50%).", False, SEC, 1),
        ("2. Padronizar a planilha-fonte: nomes de área, totais anuais, lançamentos mensais completos.", True, TXT, 0),
        ("    · Definir hoje o que significa \"(JD)\" e criar campo POLO/LOCAL.", False, SEC, 1),
        ("3. Definir metas de frequência por período (respeitando férias e sazonalidade).", True, TXT, 0),
        ("4. Acompanhamento mensal com o painel — reavaliar o mesmo indicador.", True, TXT, 0),
    ], size=16, gap=14)
    text(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.7),
         [([("Dado vira informação, informação vira decisão.", 20, True, RED)], 0, PP_ALIGN.LEFT)])


# --------------------------------------------------------------------------
def main():
    for fn in [s_capa, s_objetivo, s_fonte, s_kpis, s_sazonalidade, s_volume,
               s_distribuicao, s_area, s_atencoes, s_automotiva, s_destaques,
               s_heatmap, s_qualidade, s_plano]:
        fn()
    prs.save(OUT_PPTX)
    print("OK", OUT_PPTX)


if __name__ == "__main__":
    main()